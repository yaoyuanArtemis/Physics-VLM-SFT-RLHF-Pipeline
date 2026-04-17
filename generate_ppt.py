#!/usr/bin/env python3
"""Generate academic defense presentation for MSc thesis."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# === Colors (academic navy/blue theme) ===
NAVY = RGBColor(0x00, 0x2B, 0x5C)       # HKU-ish dark navy
ACCENT_BLUE = RGBColor(0x00, 0x6D, 0xAA) # mid blue for accents
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)  # light bg
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)

FIG_DIR = "/Users/sh01679ml/Desktop/Physics-VLM-SFT-RLHF-Pipeline/Overleaf/Figures"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, items, left=Inches(0.8), top=Inches(1.8),
                     width=Inches(11.5), font_size=20, color=DARK_GRAY, spacing=Pt(12)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def section_header(slide, title, subtitle=""):
    add_bg(slide, WHITE)
    # Navy bar at top
    add_rect(slide, 0, 0, W, Inches(0.08), NAVY)
    # Bottom bar
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), NAVY)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                title, font_size=32, color=NAVY, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
                    subtitle, font_size=18, color=MID_GRAY)
    # Thin accent line under title
    add_rect(slide, Inches(0.8), Inches(1.15 if not subtitle else 1.55),
             Inches(2), Inches(0.04), ACCENT_BLUE)


def add_page_number(slide, num, total):
    add_textbox(slide, W - Inches(1.2), H - Inches(0.55), Inches(1), Inches(0.4),
                f"{num}/{total}", font_size=11, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
        return True
    return False


# ============================================================================
# SLIDE 1: Title
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, WHITE)
add_rect(slide, 0, 0, W, Inches(3.2), NAVY)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(2.2),
            "Large Language Model Based Multi-Agent\nMaterials Scientist for Advanced Materials Design",
            font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(0.5),
            "MSc Thesis Defense",
            font_size=20, color=RGBColor(0xAA, 0xCC, 0xEE), alignment=PP_ALIGN.LEFT)

add_textbox(slide, Inches(0.8), Inches(3.8), Inches(5), Inches(2.5),
            "Liu Feng\nDepartment of Mechanical Engineering\nThe University of Hong Kong\n\nSupervisor: Prof. Wen Tongqi",
            font_size=20, color=DARK_GRAY)

add_textbox(slide, Inches(0.8), Inches(6.5), Inches(5), Inches(0.5),
            "April 2026", font_size=16, color=MID_GRAY)

TOTAL = 24  # rough estimate, will adjust

# ============================================================================
# SLIDE 2: Outline
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Outline")

outline_items = [
    "1.  Background & Motivation",
    "2.  Methodology Overview",
    "3.  Stage 1: Supervised Fine-Tuning (SFT)",
    "4.  Stage 2: RLHF with Physics-Aware Rewards",
    "5.  Stage 3: Multi-Agent Integration",
    "6.  Results & Discussion",
    "7.  Limitations & Future Work",
    "8.  Conclusion",
]
add_bullet_slide(slide, outline_items, top=Inches(2.0), font_size=24, spacing=Pt(18))
add_page_number(slide, 2, TOTAL)

# ============================================================================
# SLIDE 3: Background — The Problem
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Background", "Why do we need AI for materials characterization?")

items = [
    "•  Materials characterization is critical: quality control, failure prediction, new material design",
    "•  Current bottleneck: manual expert analysis is slow, subjective, and cannot scale",
    "•  CNNs can classify defects, but produce labels without explanation",
    '•  General-purpose VLMs fail on materials: e.g., misidentifying "dislocation networks" as "wire patterns"',
    "•  Gap: no VLM specifically adapted for materials microstructure characterization",
]
add_bullet_slide(slide, items, top=Inches(2.0), font_size=20, spacing=Pt(14))
add_page_number(slide, 3, TOTAL)

# ============================================================================
# SLIDE 4: Three Challenges
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Three Key Challenges")

# Challenge boxes
box_w = Inches(3.6)
box_h = Inches(3.8)
gap = Inches(0.5)
start_x = Inches(0.8)
top_y = Inches(2.0)

challenges = [
    ("Limited Hardware", "Single RTX 4090 (24GB)\nvs. typical 8×A100 clusters\n\n→ Need aggressive memory\n   optimization to fit 7B\n   parameter VLM training", RED),
    ("Catastrophic Forgetting", "SFT on materials data\ndestroys general capabilities\n\nChart accuracy: 86.5% → 41.2%\nPhysics violations: 12% → 34%\n\n→ Need recovery mechanism", ORANGE),
    ("Physical Consistency", "Model generates fluent but\nphysically wrong descriptions\n\nConfuses BCC/HCP phases,\nmixes nano/macro scales\n\n→ Need physics constraints\n   in training", ACCENT_BLUE),
]

for i, (title, desc, color) in enumerate(challenges):
    x = start_x + i * (box_w + gap)
    box = add_rect(slide, x, top_y, box_w, box_h, WHITE, color)
    # colored top bar
    add_rect(slide, x, top_y, box_w, Inches(0.06), color)
    add_textbox(slide, x + Inches(0.3), top_y + Inches(0.2), box_w - Inches(0.6), Inches(0.5),
                title, font_size=22, color=color, bold=True)
    add_textbox(slide, x + Inches(0.3), top_y + Inches(0.8), box_w - Inches(0.6), box_h - Inches(1.0),
                desc, font_size=16, color=DARK_GRAY)

add_page_number(slide, 4, TOTAL)

# ============================================================================
# SLIDE 5: Pipeline Overview
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Methodology: Four-Stage Pipeline")

add_image_safe(slide, f"{FIG_DIR}/pipeline_flowchart_hd-1.png",
               Inches(0.5), Inches(1.8), width=Inches(12.3))
add_page_number(slide, 5, TOTAL)

# ============================================================================
# SLIDE 6: Model & Data
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Base Model & Dataset")

items = [
    "Base Model:  Qwen2.5-VL-7B-Instruct (7 billion parameters)",
    "                    Vision encoder (ViT) + projection layer + LLM decoder",
    "",
    "Dataset:       25,263 physics material images from OmniScience",
    "                    Filtered by physics/materials keywords",
    "                    Machine-generated annotations (not expert-written)",
    "                    Covers: SEM, TEM, XRD, optical micrographs",
    "",
    "Hardware:    Single NVIDIA RTX 4090 (24GB VRAM)",
    "                    AutoDL cloud server, Ubuntu environment",
]
add_bullet_slide(slide, items, top=Inches(1.8), font_size=20, spacing=Pt(8))
add_page_number(slide, 6, TOTAL)

# ============================================================================
# SLIDE 7: VLM Architecture
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "VLM Architecture", "Qwen2.5-VL: Vision Encoder → Projection → Language Model")

add_image_safe(slide, f"{FIG_DIR}/Chapter2/vlm_architecture.png",
               Inches(1.5), Inches(1.8), width=Inches(10))
add_page_number(slide, 7, TOTAL)

# ============================================================================
# SLIDE 8: Memory Optimization
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Stage 1: Memory Optimization for Single-GPU Training")

items = [
    "Problem: Full fine-tuning of 7B model requires ~56GB — far exceeds 24GB VRAM",
    "",
    "Solution stack:",
    "  •  LoRA (rank=8, α=8) — only train q_proj & v_proj → 0.1% trainable params",
    "  •  BFloat16 mixed precision — halve memory per parameter",
    "  •  FlashAttention-2 — O(N) memory attention instead of O(N²)",
    "  •  Paged 8-bit AdamW — optimizer states in 8-bit with CPU paging",
    "  •  Gradient accumulation (16 steps) — effective batch size 16 with batch=1",
    "",
    "Result: Peak memory 23.8 GiB — fits on RTX 4090 with ~200MB to spare",
]
add_bullet_slide(slide, items, top=Inches(1.6), font_size=20, spacing=Pt(8))
add_page_number(slide, 8, TOTAL)

# ============================================================================
# SLIDE 9: SFT Training
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Stage 1: SFT Training Configuration & Results")

# Left: config
items_left = [
    "Framework: LLaMA-Factory",
    "Epochs: 3",
    "Learning rate: 5e-5 (cosine)",
    "Context length: 2048",
    "Batch size: 1 × 16 accum",
    "Warmup: 10%",
    "Training time: ~18 hours",
]
txBox = add_bullet_slide(slide, items_left, left=Inches(0.8), top=Inches(1.8),
                         width=Inches(5), font_size=18, spacing=Pt(8))

# Right: figure
add_image_safe(slide, f"{FIG_DIR}/Chapter3/sft_training_dynamics.png",
               Inches(6.5), Inches(1.6), width=Inches(6.2))
add_page_number(slide, 9, TOTAL)

# ============================================================================
# SLIDE 10: Catastrophic Forgetting
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "The Problem: Catastrophic Forgetting After SFT", "SFT improves materials knowledge but destroys general capabilities")

# Table-like layout
table_data = [
    ("Task", "Base Model", "After SFT", "Change"),
    ("Materials Defect Acc.", "68.7%", "94.3%", "▲ +25.6pp"),
    ("Chart Analysis Acc.", "86.5%", "41.2%", "▼ −45.3pp"),
    ("Physics Violations", "12%", "34%", "▼ +22pp"),
    ("Crystallographic F1", "0.63", "0.89", "▲ +0.26"),
]

top_y = Inches(2.2)
col_widths = [Inches(3.5), Inches(2.5), Inches(2.5), Inches(2.5)]
row_h = Inches(0.6)
start_x = Inches(1.2)

for r, row in enumerate(table_data):
    x = start_x
    for c, (cell, cw) in enumerate(zip(row, col_widths)):
        if r == 0:
            box = add_rect(slide, x, top_y + r * row_h, cw, row_h, NAVY)
            clr = WHITE
            bld = True
        else:
            bg = LIGHT_BLUE if r % 2 == 0 else WHITE
            box = add_rect(slide, x, top_y + r * row_h, cw, row_h, bg, RGBColor(0xCC, 0xCC, 0xCC))
            if "▼" in cell:
                clr = RED
            elif "▲" in cell:
                clr = GREEN
            else:
                clr = DARK_GRAY
            bld = False
        add_textbox(slide, x + Inches(0.15), top_y + r * row_h + Inches(0.1),
                    cw - Inches(0.3), row_h - Inches(0.2),
                    cell, font_size=18, color=clr, bold=bld,
                    alignment=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)
        x += cw

add_textbox(slide, Inches(1.2), Inches(5.5), Inches(10), Inches(0.8),
            "→ SFT teaches domain vocabulary but breaks chart reading and introduces physics errors.\n"
            "→ This motivates Stage 2: RLHF with physics-aware constraints.",
            font_size=18, color=NAVY, bold=True)
add_page_number(slide, 10, TOTAL)

# ============================================================================
# SLIDE 11: Reward Function Design
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Stage 2: Physics-Aware Reward Function", "6 components with asymmetric penalty design")

components = [
    ("Microstructural features", "+0.5 each", "grain boundary, dislocation, twin...", GREEN),
    ("Computational physics", "+0.8 each", "formation energy, stacking fault, DFT...", GREEN),
    ("Scale consistency", "−2.0", "Penalize atomistic↔macroscopic confusion", RED),
    ("Phase accuracy", "−2.5", "Penalize BCC/HCP phase reversals", RED),
    ("Non-scientific language", "−1.5", 'Penalize "artistic", "beautiful"...', RED),
    ("Ground-truth alignment", "+3.0", "Bonus when output contains reference", GREEN),
]

top_y = Inches(1.9)
for i, (name, score, desc, clr) in enumerate(components):
    y = top_y + i * Inches(0.78)
    # Score box
    score_box = add_rect(slide, Inches(0.8), y, Inches(1.5), Inches(0.6), clr)
    add_textbox(slide, Inches(0.8), y + Inches(0.1), Inches(1.5), Inches(0.4),
                score, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Name
    add_textbox(slide, Inches(2.5), y + Inches(0.1), Inches(3.5), Inches(0.4),
                name, font_size=18, color=DARK_GRAY, bold=True)
    # Description
    add_textbox(slide, Inches(6.0), y + Inches(0.1), Inches(6.5), Inches(0.4),
                desc, font_size=16, color=MID_GRAY)

add_textbox(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.5),
            "Design principle: one physics error outweighs multiple correct keywords → prioritize rigor over verbosity",
            font_size=17, color=NAVY, bold=True)
add_page_number(slide, 11, TOTAL)

# ============================================================================
# SLIDE 12: GRPO Algorithm
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Stage 2: GRPO Training", "Group Relative Policy Optimization — no critic network needed")

items = [
    "Framework: ms-swift (ModelScope)",
    "Algorithm: GRPO with ε=0.2 clipping",
    "Group size: N=4 responses per prompt",
    "Learning rate: 1e-6",
    "Context length: 1024 (reduced from 2048 for memory)",
    "Temperature: 0.8 for diverse rollouts",
    "Training: 1,811 steps, 1 epoch",
    "Peak memory: 25.7 GiB (near 24GB limit)",
]

add_bullet_slide(slide, items, left=Inches(0.8), top=Inches(1.8),
                 width=Inches(5.5), font_size=19, spacing=Pt(10))

# GRPO diagram
add_image_safe(slide, f"{FIG_DIR}/grpo.png",
               Inches(7.0), Inches(1.6), width=Inches(5.5))
add_page_number(slide, 12, TOTAL)

# ============================================================================
# SLIDE 13: GRPO Training Dynamics
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "GRPO Training Dynamics", "4 indicators over 1,811 training steps")

add_image_safe(slide, f"{FIG_DIR}/Chapter4/grpo_training_dynamics.png",
               Inches(0.5), Inches(1.6), width=Inches(12.3))

add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.7),
            "(a) Reward ↑ 3.04→3.50, std→0    (b) KL stable at 0.16    "
            "(c) Length 10→142 tokens    (d) Loss & gradient stable",
            font_size=15, color=MID_GRAY)
add_page_number(slide, 13, TOTAL)

# ============================================================================
# SLIDE 14: Full Results Table
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Results: Three-Stage Performance Comparison")

table_data = [
    ("Task", "Base Model", "SFT Only", "SFT + RLHF"),
    ("Materials Defect Acc.", "68.7%", "94.3%", "93.1%"),
    ("Chart Analysis Acc.", "86.5%", "41.2%", "78.3%"),
    ("Physics Violations ↓", "12%", "34%", "8%"),
    ("Crystallographic F1", "0.63", "0.89", "0.91"),
]

top_y = Inches(2.0)
col_widths = [Inches(3.5), Inches(2.5), Inches(2.5), Inches(2.8)]
row_h = Inches(0.7)
start_x = Inches(1.2)

for r, row in enumerate(table_data):
    x = start_x
    for c, (cell, cw) in enumerate(zip(row, col_widths)):
        if r == 0:
            box = add_rect(slide, x, top_y + r * row_h, cw, row_h, NAVY)
            clr = WHITE
            bld = True
        else:
            bg = LIGHT_BLUE if r % 2 == 0 else WHITE
            # highlight RLHF column
            if c == 3 and r > 0:
                bg = RGBColor(0xE8, 0xF5, 0xE9)
            box = add_rect(slide, x, top_y + r * row_h, cw, row_h, bg, RGBColor(0xCC, 0xCC, 0xCC))
            clr = DARK_GRAY
            bld = (c == 3)
        add_textbox(slide, x + Inches(0.15), top_y + r * row_h + Inches(0.12),
                    cw - Inches(0.3), row_h - Inches(0.2),
                    cell, font_size=20, color=clr, bold=bld,
                    alignment=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)
        x += cw

add_textbox(slide, Inches(1.2), Inches(5.8), Inches(10), Inches(1.0),
            "Key takeaways:\n"
            "•  RLHF recovers chart accuracy from 41.2% → 78.3% (residual gap: 8.2pp)\n"
            "•  Physics violations drop to 8% — below even the base model's 12%\n"
            "•  Materials accuracy barely affected: 94.3% → 93.1%",
            font_size=18, color=DARK_GRAY)
add_page_number(slide, 14, TOTAL)

# ============================================================================
# SLIDE 15: Cross-Model Comparison
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Cross-Model Comparison", "vs. Darwin1_v2 (7B scientific VLM)")

table_data = [
    ("Model", "Defect Acc.", "Chart Acc.", "Phys. Viol.↓", "Cryst. F1"),
    ("Darwin1_v2 (7B)", "75.8%", "82.3%", "15%", "0.61"),
    ("Qwen2.5-VL (Base)", "68.7%", "86.5%", "12%", "0.63"),
    ("Ours (SFT+RLHF)", "93.1%", "78.3%", "8%", "0.91"),
]

top_y = Inches(2.2)
col_widths = [Inches(3.2), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
row_h = Inches(0.7)
start_x = Inches(0.8)

for r, row in enumerate(table_data):
    x = start_x
    for c, (cell, cw) in enumerate(zip(row, col_widths)):
        if r == 0:
            box = add_rect(slide, x, top_y + r * row_h, cw, row_h, NAVY)
            clr, bld = WHITE, True
        elif r == 3:
            box = add_rect(slide, x, top_y + r * row_h, cw, row_h, RGBColor(0xE8, 0xF5, 0xE9), ACCENT_BLUE)
            clr, bld = NAVY, True
        else:
            bg = LIGHT_BLUE if r % 2 == 0 else WHITE
            box = add_rect(slide, x, top_y + r * row_h, cw, row_h, bg, RGBColor(0xCC, 0xCC, 0xCC))
            clr, bld = DARK_GRAY, False
        add_textbox(slide, x + Inches(0.1), top_y + r * row_h + Inches(0.12),
                    cw - Inches(0.2), row_h - Inches(0.2),
                    cell, font_size=19, color=clr, bold=bld,
                    alignment=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)
        x += cw

add_textbox(slide, Inches(0.8), Inches(5.3), Inches(11), Inches(1.0),
            "•  +17.3pp over Darwin on defect identification, +0.30 on crystallographic F1\n"
            "•  Lowest physics violation rate across all models\n"
            "•  Trade-off: 4pp lower chart accuracy than Darwin (specialized vs. generalist)",
            font_size=18, color=DARK_GRAY)
add_page_number(slide, 15, TOTAL)

# ============================================================================
# SLIDE 16: Qualitative Case Study
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Qualitative Comparison: Three Model Stages")

table_data = [
    ("Capability", "Base", "SFT", "RLHF"),
    ("Visual feature identification", "✓", "✓", "✓"),
    ("Domain-specific terminology", "Partial", "✓", "✓"),
    ("Structured analytical framework", "—", "✓", "✓"),
    ("Physical mechanism reasoning", "—", "—", "✓"),
    ("Property–structure linkage", "—", "—", "✓"),
    ("Instrumental limitation awareness", "—", "—", "✓"),
]

top_y = Inches(1.8)
col_widths = [Inches(5.5), Inches(2.0), Inches(2.0), Inches(2.0)]
row_h = Inches(0.6)
start_x = Inches(0.8)

for r, row in enumerate(table_data):
    x = start_x
    for c, (cell, cw) in enumerate(zip(row, col_widths)):
        if r == 0:
            box = add_rect(slide, x, top_y + r * row_h, cw, row_h, NAVY)
            clr, bld = WHITE, True
        else:
            bg = LIGHT_BLUE if r % 2 == 0 else WHITE
            box = add_rect(slide, x, top_y + r * row_h, cw, row_h, bg, RGBColor(0xCC, 0xCC, 0xCC))
            if cell == "✓":
                clr = GREEN
            elif cell == "—":
                clr = RGBColor(0xBB, 0xBB, 0xBB)
            else:
                clr = DARK_GRAY
            bld = False
        add_textbox(slide, x + Inches(0.1), top_y + r * row_h + Inches(0.08),
                    cw - Inches(0.2), row_h - Inches(0.16),
                    cell, font_size=18, color=clr, bold=bld,
                    alignment=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)
        x += cw

add_textbox(slide, Inches(0.8), Inches(6.2), Inches(10), Inches(0.5),
            "SFT adds vocabulary & structure; RLHF adds reasoning, mechanism explanation, and self-awareness of limitations",
            font_size=17, color=NAVY, bold=True)
add_page_number(slide, 16, TOTAL)

# ============================================================================
# SLIDE 17: Multi-Agent System
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Stage 3: Multi-Agent Integration", "VLM as perception core + external tools")

# Architecture description
items = [
    "Sequential pipeline architecture:",
    "",
    "  Input Image",
    "     ↓",
    "  Tool 1: VLM Perception — analyze microstructure",
    "     ↓  (extract materials, grain sizes)",
    "  Tool 2: Materials Project API — query crystal structure, band gap",
    "     ↓  (retrieve authoritative properties)",
    "  Tool 3: Physics Calculator — Hall-Petch, Bragg's law",
    "     ↓",
    "  Output: Structured JSON Report",
]
add_bullet_slide(slide, items, left=Inches(0.8), top=Inches(1.8),
                 width=Inches(5.5), font_size=18, spacing=Pt(6))

add_image_safe(slide, f"{FIG_DIR}/agent.jpg",
               Inches(7.0), Inches(1.8), height=Inches(4.8))
add_page_number(slide, 17, TOTAL)

# ============================================================================
# SLIDE 18: 12-Image Evaluation Results
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Agent Evaluation: 12-Image Quantitative Results")

eval_table = [
    ("Metric", "Value", "Notes"),
    ("VLM perception success", "12/12 (100%)", "All images received structured analysis"),
    ("Material extraction", "7/12 (58%)", "Fails on polymers & non-micrographs"),
    ("MP API query success", "8/8 (100%)", "All dispatched queries returned data"),
    ("Physics calc triggered", "4/12 (33%)", "When grain size is present"),
    ("Full 3-tool pipeline", "2/12 (17%)", "Images 158 (PbS+Pd) & 3250 (Cu)"),
    ("False positive extraction", "1/12 (8%)", '"Ti:Sapphire" → Ti'),
    ("Crash rate", "0/12 (0%)", "No runtime errors"),
]

top_y = Inches(1.8)
col_widths = [Inches(3.5), Inches(2.5), Inches(6.0)]
row_h = Inches(0.6)
start_x = Inches(0.5)

for r, row in enumerate(eval_table):
    x = start_x
    for c, (cell, cw) in enumerate(zip(row, col_widths)):
        if r == 0:
            box = add_rect(slide, x, top_y + r * row_h, cw, row_h, NAVY)
            clr, bld = WHITE, True
        else:
            bg = LIGHT_BLUE if r % 2 == 0 else WHITE
            # Highlight full pipeline row
            if r == 5:
                bg = RGBColor(0xFD, 0xE8, 0xD0)
            box = add_rect(slide, x, top_y + r * row_h, cw, row_h, bg, RGBColor(0xCC, 0xCC, 0xCC))
            clr = DARK_GRAY
            bld = (c == 1)
        add_textbox(slide, x + Inches(0.15), top_y + r * row_h + Inches(0.08),
                    cw - Inches(0.3), row_h - Inches(0.16),
                    cell, font_size=17, color=clr, bold=bld,
                    alignment=PP_ALIGN.CENTER if c == 1 else PP_ALIGN.LEFT)
        x += cw

add_page_number(slide, 18, TOTAL)

# ============================================================================
# SLIDE 18b: Key Findings from Evaluation
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Agent Key Findings", "The bottleneck is the orchestrator, not the VLM")

items = [
    "1. VLM is the reliable component (100% success)",
    "    SFT+RLHF model generalizes across nanocrystals, thin films, polymers, spectral plots",
    "",
    "2. Regex-based extraction is the main bottleneck (58%)",
    "    • Polymers (PEO, PFS-b-P2VP): not in inorganic formula list — by design",
    "    • Complex compounds: FeSe₀.₅Te₀.₅ → only 'Fe' extracted",
    '    • False positive: "Ti:Sapphire" → Ti (laser component, not a material)',
    "",
    "3. Materials Project polymorph mismatch",
    "    API returns metastable phases: PbS → monoclinic (not cubic galena),",
    "    Fe → triclinic P1 (not BCC). Needs polymorph selection logic.",
    "",
    "4. Clear improvement path",
    "    Replace regex → LLM-based structured extraction",
    "    Add polymer-specific tool modules",
    "    Implement two-pass prompting for unknown materials",
]
add_bullet_slide(slide, items, top=Inches(1.6), font_size=18, spacing=Pt(4))
add_page_number(slide, 19, TOTAL)

# ============================================================================
# SLIDE 20: Limitations
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Limitations & Honest Assessment")

items = [
    "1. Forgetting is not fully solved",
    "    Chart accuracy: 78.3% vs base 86.5% — 8.2pp residual gap remains",
    "    RLHF dataset contains only materials images, no chart-reading data",
    "",
    "2. Keyword-matching reward has blind spots",
    "    Cannot verify spatial relationships or logical coherence of reasoning",
    "    Failure case: model confused SEM/TEM scales in multi-panel figure",
    "",
    "3. Evaluation is the weakest link",
    "    Both reward function and evaluation metrics use keyword matching",
    "    Risk of circular validation — no independent human expert evaluation",
    "",
    "4. Hardware constraints shaped design, not optimality",
    "    Context length 1024, group size 4, LoRA on 2 modules only",
]
add_bullet_slide(slide, items, top=Inches(1.6), font_size=18, spacing=Pt(4))
add_page_number(slide, 20, TOTAL)

# ============================================================================
# SLIDE 21: Future Work
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Future Directions")

items = [
    "•  Better evaluation: LLM-as-Judge (GPT-4o scoring 5 dimensions) + small human expert study",
    "",
    "•  Automated reward generation: extract validation rules from textbooks/databases",
    "",
    "•  Mixed-data RLHF: interleave materials + general data to reduce forgetting",
    "",
    "•  Chain-of-thought prompting: structured per-panel analysis for multi-panel figures",
    "",
    "•  Simulation tool integration: connect VLM perception → LAMMPS/VASP → inverse design",
    "",
    "•  Cross-scale fusion: combine SEM/TEM images + XRD spectra + CIF structures",
]
add_bullet_slide(slide, items, top=Inches(1.8), font_size=20, spacing=Pt(6))
add_page_number(slide, 21, TOTAL)

# ============================================================================
# SLIDE 22: Open Source
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Open-Source Release")

table_data = [
    ("Resource", "Location"),
    ("Code & Pipeline", "github.com/yaoyuanArtemis/Physics-VLM-SFT-RLHF-Pipeline"),
    ("RLHF Model", "huggingface.co/yaoyuanlf/Qwen2.5-VL-7B-Physics-RLHF"),
    ("SFT LoRA Adapter", "huggingface.co/yaoyuanlf/qwen2.5-vl-physics-lora"),
    ("Dataset (25k)", "huggingface.co/datasets/yaoyuanlf/physics-vlm-dataset"),
]

top_y = Inches(2.2)
col_widths = [Inches(3.0), Inches(8.5)]
row_h = Inches(0.7)
start_x = Inches(0.8)

for r, row in enumerate(table_data):
    x = start_x
    for c, (cell, cw) in enumerate(zip(row, col_widths)):
        if r == 0:
            box = add_rect(slide, x, top_y + r * row_h, cw, row_h, NAVY)
            clr, bld = WHITE, True
        else:
            bg = LIGHT_BLUE if r % 2 == 0 else WHITE
            box = add_rect(slide, x, top_y + r * row_h, cw, row_h, bg, RGBColor(0xCC, 0xCC, 0xCC))
            clr, bld = DARK_GRAY if c == 0 else ACCENT_BLUE, (c == 0)
        add_textbox(slide, x + Inches(0.2), top_y + r * row_h + Inches(0.12),
                    cw - Inches(0.4), row_h - Inches(0.2),
                    cell, font_size=19, color=clr, bold=bld,
                    alignment=PP_ALIGN.LEFT)
        x += cw

add_textbox(slide, Inches(0.8), Inches(5.8), Inches(10), Inches(0.8),
            "All code, configs, reward functions, model weights, and agent demo are publicly available.\n"
            "Goal: lower the engineering barrier so others can replicate and extend this work.",
            font_size=18, color=DARK_GRAY)
add_page_number(slide, 22, TOTAL)

# ============================================================================
# SLIDE 23: Conclusion
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_header(slide, "Conclusion")

items = [
    "1.  Demonstrated that 7B VLM fine-tuning is feasible on a single 24GB GPU",
    "     through LoRA + BF16 + FlashAttention-2 + paged optimizer",
    "",
    "2.  Two-stage SFT → RLHF pipeline achieves 93.1% defect accuracy",
    "     while reducing physics violations to 8% (below base model's 12%)",
    "",
    "3.  Hand-designed physics reward function is simple but effective:",
    "     keyword-based, asymmetric penalties, empirically tuned weights",
    "",
    "4.  Multi-agent evaluation (12 images): VLM 100% reliable, but full pipeline",
    "     only 17% — bottleneck is regex orchestrator, not the model",
    "",
    "5.  Full pipeline open-sourced for reproducibility",
    "",
    "The approach is not perfect — the evaluation is limited, forgetting is not fully",
    "solved, and many design choices were memory-driven — but it works, and it runs",
    "on hardware that most academic labs already have.",
]
add_bullet_slide(slide, items, top=Inches(1.6), font_size=19, spacing=Pt(4))
add_page_number(slide, 23, TOTAL)

# ============================================================================
# SLIDE 24: Thank You
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, W, H, NAVY)

add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.5),
            "Thank You",
            font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
            "Questions & Discussion",
            font_size=28, color=RGBColor(0xAA, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(1),
            "Liu Feng  |  Department of Mechanical Engineering  |  The University of Hong Kong\n"
            "github.com/yaoyuanArtemis/Physics-VLM-SFT-RLHF-Pipeline",
            font_size=18, color=RGBColor(0x88, 0xAA, 0xCC), alignment=PP_ALIGN.CENTER)

# ============================================================================
# Update total page numbers
# ============================================================================
TOTAL = len(prs.slides)

# Save
output_path = "/Users/sh01679ml/Desktop/Physics-VLM-SFT-RLHF-Pipeline/thesis_defense.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Total slides: {TOTAL}")
